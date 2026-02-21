"""
Document Intelligence – Text Extraction Service
Supports: PDF (pdfplumber primary, PyMuPDF fallback), PPTX, PPT, TXT
"""

import io
import logging
from dataclasses import dataclass
from typing import Tuple

logger = logging.getLogger(__name__)


# ─── Result container ────────────────────────────────────────────────────────

@dataclass
class ExtractionResult:
    text: str           # Full structured extracted text
    file_type: str      # pdf | pptx | ppt | txt
    page_count: int     # Total pages / slides
    scanned: bool       # True if image-only / no selectable text
    warning: str        # Human-readable warning, empty if clean


# ─── Public entry point ──────────────────────────────────────────────────────

def extract_text_from_file(content: bytes, filename: str) -> Tuple[str, str]:
    """
    Legacy-compatible wrapper.
    Returns (extracted_text, file_type).
    Raises ValueError with a user-friendly message on failure.
    """
    if not content:
        raise ValueError("The uploaded file is empty. Please upload a file with content.")

    result = extract_document(content, filename)

    # Surface warnings as part of the text so the frontend can display them
    if result.scanned:
        return result.text, result.file_type
    if result.warning:
        full = f"[⚠️ {result.warning}]\n\n{result.text}".strip()
        return full, result.file_type

    return result.text, result.file_type


def extract_document(content: bytes, filename: str) -> ExtractionResult:
    """
    Full extraction returning a rich ExtractionResult.
    """
    if not content:
        raise ValueError("The uploaded file is empty. Please upload a file with content.")

    lower = filename.lower()

    if lower.endswith(".pdf"):
        return _extract_pdf(content)
    elif lower.endswith(".pptx"):
        return _extract_pptx(content)
    elif lower.endswith(".ppt"):
        return _extract_ppt(content)
    elif lower.endswith(".txt"):
        return _extract_txt(content)
    else:
        ext = filename.rsplit(".", 1)[-1].upper() if "." in filename else "unknown"
        raise ValueError(
            f"Unsupported file type: .{ext}. "
            "Please upload a PDF, PowerPoint (.pptx / .ppt), or plain text file."
        )


# ─── PDF extraction ──────────────────────────────────────────────────────────

def _extract_pdf(content: bytes) -> ExtractionResult:
    """
    Primary: pdfplumber  →  Fallback: PyMuPDF (fitz)
    Detects scanned / image-only PDFs.
    """
    # ── pdfplumber (primary) ──────────────────────────────────────
    try:
        import pdfplumber

        parts: list[str] = []
        scanned_pages: list[int] = []
        total = 0

        with pdfplumber.open(io.BytesIO(content)) as pdf:
            total = len(pdf.pages)

            if total == 0:
                raise ValueError("This PDF has no pages.")

            for i, page in enumerate(pdf.pages, 1):
                raw = page.extract_text()
                text = (raw or "").strip()

                if text:
                    parts.append(f"Page {i}:\n{text}")
                else:
                    scanned_pages.append(i)
                    parts.append(
                        f"Page {i}:\n[No extractable text — this page may be an image or scan]"
                    )

        all_scanned = len(scanned_pages) == total
        some_scanned = bool(scanned_pages) and not all_scanned
        extracted = "\n\n".join(parts)

        if all_scanned:
            return ExtractionResult(
                text=(
                    "This appears to be a scanned document. OCR not enabled.\n\n"
                    "Please copy and paste your text manually into the editor below, "
                    "or use a PDF with selectable text."
                ),
                file_type="pdf",
                page_count=total,
                scanned=True,
                warning="",
            )

        warning = (
            f"Pages {', '.join(str(p) for p in scanned_pages)} appear to be scanned images "
            "and could not be extracted."
            if some_scanned
            else ""
        )

        return ExtractionResult(
            text=extracted,
            file_type="pdf",
            page_count=total,
            scanned=False,
            warning=warning,
        )

    except ValueError:
        raise
    except Exception as plumber_err:
        logger.warning("pdfplumber failed (%s), trying PyMuPDF fallback", plumber_err)

    # ── PyMuPDF fallback ─────────────────────────────────────────
    try:
        import fitz  # PyMuPDF

        parts: list[str] = []
        scanned_pages: list[int] = []

        doc = fitz.open(stream=content, filetype="pdf")
        total = doc.page_count

        if total == 0:
            raise ValueError("This PDF has no pages.")

        for i in range(total):
            page = doc[i]
            text = page.get_text("text").strip()  # type: ignore[attr-defined]
            page_num = i + 1

            if text:
                parts.append(f"Page {page_num}:\n{text}")
            else:
                scanned_pages.append(page_num)
                parts.append(
                    f"Page {page_num}:\n[No extractable text — this page may be an image or scan]"
                )

        doc.close()

        all_scanned = len(scanned_pages) == total
        some_scanned = bool(scanned_pages) and not all_scanned
        extracted = "\n\n".join(parts)

        if all_scanned:
            return ExtractionResult(
                text=(
                    "This appears to be a scanned document. OCR not enabled.\n\n"
                    "Please copy and paste your text manually into the editor below, "
                    "or use a PDF with selectable text."
                ),
                file_type="pdf",
                page_count=total,
                scanned=True,
                warning="",
            )

        warning = (
            f"Pages {', '.join(str(p) for p in scanned_pages)} appear to be scanned images "
            "and could not be extracted."
            if some_scanned
            else ""
        )

        return ExtractionResult(
            text=extracted,
            file_type="pdf",
            page_count=total,
            scanned=False,
            warning=warning,
        )

    except ValueError:
        raise
    except Exception as fitz_err:
        logger.error("PyMuPDF also failed: %s", fitz_err)
        raise ValueError(
            "Could not read this PDF. The file may be corrupted, password-protected, or in an "
            "unsupported format. Please try a different file."
        )


# ─── PPTX extraction ─────────────────────────────────────────────────────────

def _extract_pptx(content: bytes) -> ExtractionResult:
    try:
        from pptx import Presentation
        from pptx.util import Pt  # noqa: F401 – ensure pptx is importable

        prs = Presentation(io.BytesIO(content))
        slides = prs.slides
        total = len(slides)

        if total == 0:
            raise ValueError("This PowerPoint file has no slides.")

        parts: list[str] = []
        empty_slides: list[int] = []

        for i, slide in enumerate(slides, 1):
            texts: list[str] = []

            for shape in slide.shapes:
                # Text frames (normal text boxes, titles, content placeholders)
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs).strip()
                        if line:
                            texts.append(line)

                # Tables
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            texts.append(" | ".join(cells))

            if texts:
                parts.append(f"Slide {i}:\n" + "\n".join(texts))
            else:
                empty_slides.append(i)
                parts.append(f"Slide {i}:\n[No text content on this slide]")

        extracted = "\n\n".join(parts)
        warning = (
            f"Slides {', '.join(str(s) for s in empty_slides)} had no text content."
            if empty_slides
            else ""
        )

        return ExtractionResult(
            text=extracted,
            file_type="pptx",
            page_count=total,
            scanned=False,
            warning=warning,
        )

    except ValueError:
        raise
    except Exception as e:
        raise ValueError(
            f"Could not read this PowerPoint file: {e}. "
            "Please ensure it is a valid .pptx file."
        )


def _extract_ppt(content: bytes) -> ExtractionResult:
    """
    Old binary .ppt format — attempt as pptx first (python-pptx can sometimes open them),
    then give a clear user-facing error.
    """
    try:
        result = _extract_pptx(content)
        result.file_type = "ppt"
        return result
    except Exception:
        raise ValueError(
            "Old-format .ppt files are not fully supported. "
            "Please re-save your presentation as .pptx (PowerPoint 2007+) and upload again."
        )


# ─── Plain text ──────────────────────────────────────────────────────────────

def _extract_txt(content: bytes) -> ExtractionResult:
    try:
        text = content.decode("utf-8", errors="replace").strip()
    except Exception:
        text = content.decode("latin-1", errors="replace").strip()

    if not text:
        raise ValueError("This text file appears to be empty.")

    return ExtractionResult(
        text=text,
        file_type="txt",
        page_count=1,
        scanned=False,
        warning="",
    )
